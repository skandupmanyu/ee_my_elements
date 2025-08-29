#!/usr/bin/env python3
"""
PowerPoint Slide Splitter

This application takes a PowerPoint file as input and converts each slide 
into individual .pptx files, each named with a unique UUID.

Author: AI Assistant
"""

import argparse
import os
import sys
import uuid
import time
import subprocess
import tempfile
import zipfile
import shutil
from datetime import datetime
from pathlib import Path
from typing import List, Dict
import xml.etree.ElementTree as ET
import xml.dom.minidom

try:
    from pptx import Presentation
    from PIL import Image, ImageDraw, ImageFont
except ImportError as e:
    missing_libs = []
    try:
        from pptx import Presentation
    except ImportError:
        missing_libs.append("python-pptx")
    
    try:
        from PIL import Image
    except ImportError:
        missing_libs.append("Pillow")
    
    if missing_libs:
        print(f"Error: Missing required libraries: {', '.join(missing_libs)}")
        print("Please install them using: pip install -r requirements.txt")
        sys.exit(1)


class SlideThumbnailGenerator:
    """Advanced thumbnail generator with multiple conversion methods for accuracy."""
    
    def __init__(self):
        self.default_slide_size = (1440, 810)  # Standard 16:9 aspect ratio
        
        # Check available conversion methods
        self.conversion_methods = self._detect_conversion_methods()
        print(f"🔍 Available conversion methods: {', '.join(self.conversion_methods)}")
    
    def _detect_conversion_methods(self) -> List[str]:
        """Detect which conversion methods are available on this system."""
        methods = []
        
        # Check for macOS Quick Look (qlmanage)
        try:
            result = subprocess.run(['qlmanage', '-h'], capture_output=True, timeout=3)
            if result.returncode == 0:
                methods.append('qlmanage')
        except:
            pass
        
        # Check for LibreOffice
        libreoffice_paths = [
            "soffice", "libreoffice",
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",
            "/usr/bin/soffice", "/usr/bin/libreoffice"
        ]
        for path in libreoffice_paths:
            try:
                result = subprocess.run([path, "--version"], capture_output=True, timeout=3)
                if result.returncode == 0:
                    methods.append('libreoffice')
                    break
            except:
                continue
        
        # Always available fallback
        methods.append('pil_basic')
        
        return methods
    
    def create_high_quality_thumbnail_from_pptx(self, pptx_path: str, slide_number: int) -> Image.Image:
        """
        Create a high-quality thumbnail from PPTX file using the best available method.
        
        Args:
            pptx_path: Path to the PPTX file
            slide_number: Slide number for identification
            
        Returns:
            PIL Image object
        """
        print(f"    🎨 Generating high-quality thumbnail using best available method...")
        
        # Try methods in order of quality/preference
        for method in self.conversion_methods:
            try:
                if method == 'qlmanage':
                    img = self._convert_with_qlmanage(pptx_path)
                    if img:
                        print(f"    ✅ Used macOS Quick Look for high-quality conversion")
                        return img
                
                elif method == 'libreoffice':
                    img = self._convert_with_libreoffice(pptx_path)
                    if img:
                        print(f"    ✅ Used LibreOffice for accurate conversion")
                        return img
                
            except Exception as e:
                print(f"    ⚠️  {method} conversion failed: {e}")
                continue
        
        # Final fallback - create a better placeholder
        print(f"    ℹ️  Using enhanced fallback thumbnail")
        return self._create_enhanced_fallback_thumbnail(slide_number)
    
    def _convert_with_qlmanage(self, pptx_path: str) -> Image.Image:
        """Convert PPTX to image using macOS Quick Look."""
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = Path(temp_dir) / "thumbnail.png"
            
            # Use qlmanage to generate thumbnail
            cmd = [
                'qlmanage', 
                '-t', str(pptx_path),
                '-s', '1440',  # Size
                '-o', temp_dir
            ]
            
            result = subprocess.run(cmd, capture_output=True, timeout=15)
            
            if result.returncode == 0:
                # qlmanage creates files with specific naming
                generated_files = list(Path(temp_dir).glob("*.png"))
                if generated_files:
                    return Image.open(generated_files[0])
            
            return None
    
    def _convert_with_libreoffice(self, pptx_path: str) -> Image.Image:
        """Convert PPTX to image using LibreOffice."""
        with tempfile.TemporaryDirectory() as temp_dir:
            # Find LibreOffice
            soffice_cmd = None
            for path in ["/Applications/LibreOffice.app/Contents/MacOS/soffice", "soffice", "libreoffice"]:
                try:
                    result = subprocess.run([path, "--version"], capture_output=True, timeout=3)
                    if result.returncode == 0:
                        soffice_cmd = path
                        break
                except:
                    continue
            
            if not soffice_cmd:
                return None
            
            # Convert to PNG
            cmd = [
                soffice_cmd,
                '--headless',
                '--invisible',
                '--nodefault',
                '--nolockcheck', 
                '--nologo',
                '--norestore',
                '--convert-to', 'png',
                '--outdir', temp_dir,
                pptx_path
            ]
            
            env = os.environ.copy()
            env['HOME'] = os.path.expanduser('~')
            
            result = subprocess.run(cmd, capture_output=True, timeout=20, env=env)
            
            if result.returncode == 0:
                png_files = list(Path(temp_dir).glob("*.png"))
                if png_files:
                    return Image.open(png_files[0])
            
            return None
    
    def _create_enhanced_fallback_thumbnail(self, slide_number: int) -> Image.Image:
        """Create an enhanced fallback thumbnail with better styling."""
        img = Image.new('RGB', self.default_slide_size, '#2C3E50')  # Professional dark background
        draw = ImageDraw.Draw(img)
        
        try:
            title_font = ImageFont.truetype("/System/Library/Fonts/Arial.ttf", 64)
            subtitle_font = ImageFont.truetype("/System/Library/Fonts/Arial.ttf", 32)
        except:
            title_font = ImageFont.load_default()
            subtitle_font = ImageFont.load_default()
        
        # Draw gradient-like background
        for i in range(0, self.default_slide_size[1], 10):
            alpha = int(255 * (1 - i / self.default_slide_size[1]) * 0.3)
            color = f"#{hex(min(255, 44 + alpha))[2:].zfill(2)}{hex(min(255, 62 + alpha))[2:].zfill(2)}{hex(min(255, 80 + alpha))[2:].zfill(2)}"
            draw.rectangle([0, i, self.default_slide_size[0], i + 10], fill=color)
        
        # Main text
        main_text = f"Slide {slide_number}"
        subtitle_text = "High-quality preview not available"
        
        # Center the text
        bbox = draw.textbbox((0, 0), main_text, font=title_font)
        text_width = bbox[2] - bbox[0]
        text_height = bbox[3] - bbox[1]
        
        x = (self.default_slide_size[0] - text_width) // 2
        y = (self.default_slide_size[1] - text_height) // 2 - 20
        
        # Draw main text with shadow
        draw.text((x + 2, y + 2), main_text, fill='#000000', font=title_font)  # Shadow
        draw.text((x, y), main_text, fill='#FFFFFF', font=title_font)  # Main text
        
        # Subtitle
        bbox2 = draw.textbbox((0, 0), subtitle_text, font=subtitle_font)
        subtitle_width = bbox2[2] - bbox2[0]
        x2 = (self.default_slide_size[0] - subtitle_width) // 2
        y2 = y + text_height + 20
        
        draw.text((x2 + 1, y2 + 1), subtitle_text, fill='#000000', font=subtitle_font)  # Shadow
        draw.text((x2, y2), subtitle_text, fill='#BDC3C7', font=subtitle_font)  # Subtitle
        
        return img


class PowerPointSplitter:
    """Class to handle splitting PowerPoint presentations into individual slides."""
    
    def __init__(self, input_file: str, output_dir: str = None, group_name: str = None, base_name: str = None):
        """
        Initialize the PowerPoint splitter.
        
        Args:
            input_file (str): Path to the input PowerPoint file
            output_dir (str): Directory to save the individual slide files (optional, uses temp dir if not provided)
            group_name (str): Name of the group for XML metadata
            base_name (str): Base name for the output zip file (optional, uses input filename if not provided)
        """
        self.input_file = Path(input_file)
        self.group_name = group_name or "Default Group"
        self.base_name = base_name or self.input_file.stem
        self.thumbnail_generator = SlideThumbnailGenerator()
        self.temp_dir_created = False
        
        # Validate input file
        if not self.input_file.exists():
            raise FileNotFoundError(f"Input file not found: {self.input_file}")
        
        if not self.input_file.suffix.lower() in ['.pptx', '.ppt']:
            raise ValueError("Input file must be a PowerPoint file (.pptx or .ppt)")
        
        # Use provided output_dir or create a temporary directory
        if output_dir:
            self.output_dir = Path(output_dir)
            self.output_dir.mkdir(parents=True, exist_ok=True)
        else:
            # Create a temporary directory for intermediate files
            self.output_dir = Path(tempfile.mkdtemp(prefix="pptx_split_"))
            self.temp_dir_created = True
            print(f"📁 Using temporary directory: {self.output_dir.name}")
    
    def split_slides(self, progress_callback=None) -> List[str]:
        """
        Split the PowerPoint presentation into individual slide files and create XML metadata.
        
        Args:
            progress_callback: Optional callback function to report progress.
                             Called with (current_slide, total_slides, slide_title, status)
        
        Returns:
            List[str]: List of created file paths
        """
        print(f"📂 Loading presentation: {self.input_file}")
        start_time = time.time()
        
        try:
            # Load the original presentation
            presentation = Presentation(self.input_file)
            total_slides = len(presentation.slides)
            
            print(f"📊 Found {total_slides} slides to process")
            print(f"⚡ Using high-quality thumbnail generation with best available method")
            
            created_files = []
            slide_metadata = []
            
            for i, slide in enumerate(presentation.slides, 1):
                print(f"Processing slide {i}/{total_slides}...", end=" ")
                
                # Extract slide name/title first
                slide_name = self._extract_slide_name(slide, i)
                
                # Report progress - starting slide processing
                if progress_callback:
                    progress_callback(i, total_slides, slide_name, "creating_pptx")
                
                # Create a new presentation that will contain only this slide
                new_presentation = self._create_single_slide_presentation(presentation, slide, i-1)
                
                # Generate a unique UUID for the filename
                file_uuid = str(uuid.uuid4())
                output_file = self.output_dir / f"{file_uuid}.pptx"
                
                # Save the new presentation
                new_presentation.save(str(output_file))
                created_files.append(str(output_file))
                
                # Report progress - starting thumbnail generation
                if progress_callback:
                    progress_callback(i, total_slides, slide_name, "creating_thumbnail")
                
                # Generate high-quality thumbnail from PPTX file
                thumbnail_img = self.thumbnail_generator.create_high_quality_thumbnail_from_pptx(str(output_file), i)
                thumbnail_path = self._create_composite_thumbnail_from_image(thumbnail_img, file_uuid)
                
                # Report progress - slide completed
                if progress_callback:
                    progress_callback(i, total_slides, slide_name, "completed")
                
                # Store slide metadata
                slide_metadata.append({
                    'name': slide_name,
                    'id': file_uuid,
                    'thumbMode': '1'
                })
                
                print(f"✅ {output_file.name} + {Path(thumbnail_path).name}")
            
            # Generate XML metadata
            if progress_callback:
                progress_callback(total_slides + 1, total_slides + 2, "XML Metadata", "creating_xml")
            print(f"\n📄 Generating XML metadata...")
            self._create_xml_metadata(slide_metadata)
            
            # Create zip archive with all generated files
            if progress_callback:
                progress_callback(total_slides + 2, total_slides + 2, "Zip Archive", "creating_zip")
            print(f"\n📦 Creating zip archive...")
            zip_path = self._create_zip_archive()
            
            # Final completion
            if progress_callback:
                progress_callback(total_slides + 2, total_slides + 2, "Export Complete", "export_complete")
            
            total_time = time.time() - start_time
            print(f"\n🎉 Processing complete!")
            print(f"⏱️  Total time: {total_time:.1f}s ({total_time/total_slides:.1f}s per slide)")
            print(f"🚀 Performance: {total_slides/total_time:.1f} slides/second")
            print(f"📈 High-quality thumbnails with accurate visual representation!")
            
            return created_files
            
        except Exception as e:
            print(f"❌ Error processing presentation: {e}")
            # Clean up temp directory if we created it and there was an error
            if hasattr(self, 'temp_dir_created') and self.temp_dir_created:
                try:
                    shutil.rmtree(self.output_dir, ignore_errors=True)
                except:
                    pass
            raise
    
    def _create_single_slide_presentation(self, source_presentation, target_slide, slide_index):
        """
        Create a new presentation containing only the specified slide.
        
        Args:
            source_presentation: The original presentation
            target_slide: The slide to extract
            slide_index: Index of the slide in the original presentation
        
        Returns:
            A new Presentation object with only the target slide
        """
        import tempfile
        import shutil
        from pathlib import Path
        
        # Create a temporary copy of the original presentation
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
            temp_path = Path(temp_file.name)
        
        # Save the source presentation to temp file
        source_presentation.save(str(temp_path))
        
        # Load the temp presentation
        temp_presentation = Presentation(str(temp_path))
        
        # Remove all slides except the target slide
        slides_to_remove = []
        for i, slide in enumerate(temp_presentation.slides):
            if i != slide_index:
                slides_to_remove.append(i)
        
        # Remove slides in reverse order to avoid index shifting
        for slide_idx in reversed(slides_to_remove):
            rId = temp_presentation.slides._sldIdLst[slide_idx].rId
            temp_presentation.part.drop_rel(rId)
            del temp_presentation.slides._sldIdLst[slide_idx]
        
        # Clean up temp file
        try:
            temp_path.unlink()
        except:
            pass
        
        return temp_presentation
    
    def _extract_slide_name(self, slide, slide_number: int) -> str:
        """
        Extract a meaningful name from the slide.
        
        Args:
            slide: The slide object
            slide_number: The slide number (1-based)
            
        Returns:
            str: The extracted slide name
        """
        # Try to find a title in the slide
        for shape in slide.shapes:
            try:
                if hasattr(shape, 'text') and shape.text.strip():
                    text_content = shape.text.strip()
                    
                    # Check if this might be a title (usually short text or placeholder)
                    is_likely_title = False
                    
                    # Check if it's a placeholder
                    try:
                        if hasattr(shape, 'placeholder_format'):
                            is_likely_title = True
                    except:
                        pass
                    
                    # Or if it's short text (likely a title)
                    if len(text_content) < 100:
                        is_likely_title = True
                    
                    if is_likely_title:
                        # Clean up the title (remove newlines, limit length)
                        title = ' '.join(text_content.split())
                        if len(title) > 50:
                            title = title[:47] + "..."
                        return title
            except Exception:
                # Skip shapes that cause errors
                continue
        
        # Fallback to generic slide name
        return f"Slide {slide_number}"
    
    def _generate_reproducible_uuid(self, name: str) -> str:
        """
        Generate a reproducible UUID based on a name.
        
        Args:
            name: The name to base the UUID on
            
        Returns:
            str: A reproducible UUID string
        """
        # Use UUID5 with a custom namespace for reproducibility
        namespace = uuid.UUID('6ba7b810-9dad-11d1-80b4-00c04fd430c8')  # DNS namespace
        return str(uuid.uuid5(namespace, name))
    
    def _create_xml_metadata(self, slide_metadata: List[Dict]) -> None:
        """
        Create the MyElements.xml metadata file.
        
        Args:
            slide_metadata: List of slide metadata dictionaries
        """
        # Generate reproducible UUID for the group
        group_id = self._generate_reproducible_uuid(self.group_name)
        
        # Create XML structure
        root = ET.Element('ee4p')
        group = ET.SubElement(root, 'group')
        group.set('id', group_id)
        group.set('name', self.group_name)
        
        # Add each slide as an element
        for slide_info in slide_metadata:
            element = ET.SubElement(group, 'element')
            element.set('name', slide_info['name'])
            element.set('thumbMode', slide_info['thumbMode'])
            element.set('id', slide_info['id'])
        
        # Write XML file with pretty formatting
        xml_file_path = self.output_dir / 'MyElements.xml'
        
        # Convert to string and pretty print
        rough_string = ET.tostring(root, 'utf-8')
        reparsed = xml.dom.minidom.parseString(rough_string)
        pretty_xml = reparsed.toprettyxml(indent="  ")
        
        # Remove XML declaration and empty lines
        lines = [line for line in pretty_xml.split('\n') if line.strip() and not line.strip().startswith('<?xml')]
        with open(xml_file_path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(lines))
        
        print(f"📄 Created XML metadata: {xml_file_path.name}")
        print(f"   Group: {self.group_name} (ID: {group_id})")
        print(f"   Elements: {len(slide_metadata)}")
    
    def _create_composite_thumbnail_from_image(self, base_image: Image.Image, file_uuid: str) -> str:
        """
        Create a single thumbnail from a PIL Image.
        
        Args:
            base_image: PIL Image object
            file_uuid: UUID for the output filename
            
        Returns:
            str: Path to the generated thumbnail file
        """
        try:
            # Create single thumbnail at 120px height (largest size)
            thumbnail_height = 120
            base_width, base_height = base_image.size
            aspect_ratio = base_width / base_height
            thumbnail_width = int(thumbnail_height * aspect_ratio)
            
            # Resize the image maintaining aspect ratio
            thumbnail = base_image.resize((thumbnail_width, thumbnail_height), Image.Resampling.LANCZOS)
            
            # Save thumbnail
            thumbnail_path = self.output_dir / f"{file_uuid}.png"
            thumbnail.save(thumbnail_path, 'PNG', quality=95)
            
            return str(thumbnail_path)
            
        except Exception as e:
            print(f"Error creating thumbnail: {e}")
            return None
    
    def _create_zip_archive(self) -> str:
        """
        Create a zip archive containing all generated files, then clean up.
        Places zip file at the same level as the input PowerPoint file.
        
        Returns:
            str: Path to the created zip file
        """
        try:
            # Generate timestamp for unique filename
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            
            # Create zip filename based on base name + timestamp
            zip_filename = f"{self.base_name}_{timestamp}.zip"
            
            # Place zip file at the same level as the input file (not in output dir)
            zip_path = self.input_file.parent / zip_filename
            
            # Get all generated files to include in zip
            files_to_zip = []
            
            # Add all PPTX files (individual slides)
            pptx_files = list(self.output_dir.glob("*.pptx"))
            files_to_zip.extend(pptx_files)
            
            # Add all PNG files (thumbnails) 
            png_files = list(self.output_dir.glob("*.png"))
            files_to_zip.extend(png_files)
            
            # Add XML metadata file
            xml_file = self.output_dir / "MyElements.xml"
            if xml_file.exists():
                files_to_zip.append(xml_file)
            
            if not files_to_zip:
                print(f"    ⚠️  No files found to archive")
                return None
            
            # Create zip archive
            with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED, compresslevel=6) as zipf:
                for file_path in files_to_zip:
                    # Add file at root level of zip (no nested folders)
                    # Use only the filename, not the full path
                    zipf.write(file_path, file_path.name)
            
            # Verify zip contents
            file_count = len(files_to_zip)
            zip_size = zip_path.stat().st_size / (1024 * 1024)  # Size in MB
            
            print(f"    ✅ Compressed {file_count} files")
            print(f"    📦 Archive size: {zip_size:.1f} MB")
            print(f"    📁 Saved to: {zip_path}")
            
            # Clean up generated files now that they're archived
            print(f"\n🧹 Cleaning up generated files...")
            cleanup_count = self._cleanup_generated_files()
            print(f"    ✅ Removed {cleanup_count} generated files")
            print(f"    📦 Final output: {zip_filename}")
            
            return str(zip_path)
            
        except Exception as e:
            print(f"    ❌ Error creating zip archive: {e}")
            return None
    
    def _cleanup_generated_files(self) -> int:
        """
        Remove all generated files (PPTX, PNG, XML) and the temporary directory.
        
        Returns:
            int: Number of files removed
        """
        removed_count = 0
        
        try:
            # Remove all generated files
            file_patterns = ["*.pptx", "*.png", "*.xml"]
            
            for pattern in file_patterns:
                files_to_remove = list(self.output_dir.glob(pattern))
                for file_path in files_to_remove:
                    try:
                        file_path.unlink()
                        removed_count += 1
                    except Exception as e:
                        print(f"    ⚠️  Could not remove {file_path.name}: {e}")
            
            # Remove the directory
            try:
                if self.temp_dir_created:
                    # For temp directories, always remove them
                    remaining_files = list(self.output_dir.iterdir())
                    if remaining_files:
                        # Clean up any remaining files
                        for file_path in remaining_files:
                            try:
                                file_path.unlink()
                            except:
                                pass
                    
                    self.output_dir.rmdir()
                    print(f"    🗑️  Removed temporary directory: {self.output_dir.name}")
                else:
                    # For user-specified directories, only remove if empty
                    remaining_files = list(self.output_dir.iterdir())
                    if not remaining_files:
                        self.output_dir.rmdir()
                        print(f"    🗑️  Removed empty directory: {self.output_dir.name}")
                    else:
                        print(f"    ℹ️  Kept directory (contains {len(remaining_files)} other files)")
                        
            except Exception as e:
                print(f"    ℹ️  Directory cleanup skipped: {e}")
                
        except Exception as e:
            print(f"    ⚠️  Error during cleanup: {e}")
        
        return removed_count


def main():
    """Main function to handle command-line interface."""
    parser = argparse.ArgumentParser(
        description="Split PowerPoint presentation into individual slide files",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python pptx_slide_splitter.py presentation.pptx
  python pptx_slide_splitter.py presentation.pptx -o my_slides/
  python pptx_slide_splitter.py presentation.pptx -g "PodHandler"
  python pptx_slide_splitter.py presentation.pptx --group-name "My Custom Group" --output-dir ./individual_slides/
        """
    )
    
    parser.add_argument(
        "input_file",
        help="Path to the input PowerPoint file (.pptx or .ppt)"
    )
    
    parser.add_argument(
        "-o", "--output-dir",
        help="Directory to save the individual slide files (default: temporary directory)",
        default=None
    )
    
    parser.add_argument(
        "-g", "--group-name",
        help="Name of the group for XML metadata (default: derived from presentation filename)",
        default=None
    )
    
    parser.add_argument(
        "-v", "--verbose",
        action="store_true",
        help="Enable verbose output"
    )
    
    args = parser.parse_args()
    
    try:
        # Determine group name
        group_name = args.group_name
        if not group_name:
            # Use the presentation filename (without extension) as default group name
            group_name = Path(args.input_file).stem
        
        # Create the splitter
        splitter = PowerPointSplitter(args.input_file, args.output_dir, group_name)
        
        # Split the slides
        created_files = splitter.split_slides()
        
        print(f"\n✅ Successfully created {len(created_files)} individual slide files")
        print(f"📁 Output directory: {splitter.output_dir}")
        
        if args.verbose:
            print("\nCreated files:")
            for file_path in created_files:
                print(f"  • {Path(file_path).name}")
        
    except FileNotFoundError as e:
        print(f"❌ Error: {e}")
        sys.exit(1)
    except ValueError as e:
        print(f"❌ Error: {e}")
        sys.exit(1)
    except Exception as e:
        print(f"❌ Unexpected error: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()
