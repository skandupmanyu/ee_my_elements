"""
PowerPoint presentation splitter module for Export for My Efficient Elements.

This module contains the main PowerPointSplitter class responsible for:
- Splitting PowerPoint presentations into individual slide files
- Generating high-quality thumbnails
- Creating XML metadata
- Managing file operations and cleanup
"""

import tempfile
import time
from pathlib import Path
from typing import List, Dict, Optional, Callable

try:
    from pptx import Presentation
except ImportError as e:
    print("Error: Missing required library: python-pptx")
    print("Please install it using: pip install python-pptx")
    raise

from config.settings import get_processing_config, SUPPORTED_FILE_TYPES
from src.core.thumbnail_generator import SlideThumbnailGenerator
from src.core.xml_generator import XMLGenerator
from src.utils.uuid_utils import generate_unique_uuid, generate_reproducible_uuid
from src.utils.file_utils import (
    create_temp_directory, 
    create_zip_archive, 
    cleanup_files,
    cleanup_directory,
    generate_timestamped_filename,
    get_file_size_mb,
    validate_file_access
)


class PowerPointSplitter:
    """Class to handle splitting PowerPoint presentations into individual slides."""
    
    def __init__(
        self, 
        input_file: str, 
        output_dir: Optional[str] = None, 
        group_name: Optional[str] = None, 
        base_name: Optional[str] = None
    ):
        """
        Initialize the PowerPoint splitter.
        
        Args:
            input_file: Path to the input PowerPoint file
            output_dir: Directory to save the individual slide files (optional, uses temp dir if not provided)
            group_name: Name of the group for XML metadata
            base_name: Base name for the output zip file (optional, uses input filename if not provided)
        """
        self.config = get_processing_config()
        self.input_file = Path(input_file)
        self.group_name = group_name or "Default Group"
        self.base_name = base_name or self.input_file.stem
        
        # Initialize generators
        self.thumbnail_generator = SlideThumbnailGenerator()
        self.xml_generator = XMLGenerator()
        
        # Track if we created a temporary directory
        self.temp_dir_created = False
        
        # Validate input file
        self._validate_input_file()
        
        # Set up output directory
        self._setup_output_directory(output_dir)
    
    def _validate_input_file(self) -> None:
        """Validate the input PowerPoint file."""
        is_valid, error_msg = validate_file_access(self.input_file)
        if not is_valid:
            raise FileNotFoundError(error_msg)
        
        if not self.input_file.suffix.lower() in [f'.{ext}' for ext in SUPPORTED_FILE_TYPES]:
            supported_ext = ', '.join([f'.{ext}' for ext in SUPPORTED_FILE_TYPES])
            raise ValueError(f"Input file must be a PowerPoint file ({supported_ext})")
    
    def _setup_output_directory(self, output_dir: Optional[str]) -> None:
        """Set up the output directory for processing."""
        if output_dir:
            self.output_dir = Path(output_dir)
            self.output_dir.mkdir(parents=True, exist_ok=True)
        else:
            # Create a temporary directory for intermediate files
            self.output_dir = create_temp_directory(self.config['temp_prefix'])
            self.temp_dir_created = True
            if self.config['verbose']:
                print(f"📁 Using temporary directory: {self.output_dir.name}")
    
    def split_slides(self, progress_callback: Optional[Callable] = None) -> List[str]:
        """
        Split the PowerPoint presentation into individual slide files and create XML metadata.
        
        Args:
            progress_callback: Optional callback function to report progress.
                             Called with (current_slide, total_slides, slide_title, status)
        
        Returns:
            List of created file paths
        """
        if self.config['verbose']:
            print(f"📂 Loading presentation: {self.input_file}")
        
        start_time = time.time()
        
        try:
            # Load the original presentation
            presentation = Presentation(self.input_file)
            total_slides = len(presentation.slides)
            
            if self.config['verbose']:
                print(f"📊 Found {total_slides} slides to process")
                print(f"⚡ Using high-quality thumbnail generation with best available method")
            
            created_files = []
            slide_metadata = []
            
            # Optimized bulk thumbnail generation - convert entire presentation to PDF once
            if self.config['verbose']:
                print(f"🚀 Using optimized bulk thumbnail generation...")
            
            # Generate all thumbnails at once using bulk conversion
            bulk_thumbnail_paths = self.thumbnail_generator.create_high_quality_thumbnails_bulk(
                self.input_file, total_slides
            )
            
            # Process each slide with pre-generated thumbnails
            for i, slide in enumerate(presentation.slides, 1):
                if self.config['verbose']:
                    print(f"Processing slide {i}/{total_slides}...", end=" ")
                
                # Extract slide name/title first
                slide_name = self._extract_slide_name(slide, i)
                
                # Report progress - starting slide processing
                if progress_callback:
                    progress_callback(i, total_slides, slide_name, "creating_pptx")
                
                # Create individual slide presentation
                new_presentation = self._create_single_slide_presentation(presentation, slide, i-1)
                
                # Generate unique UUID for the filename
                file_uuid = generate_unique_uuid()
                output_file = self.output_dir / f"{file_uuid}.pptx"
                
                # Save the new presentation
                new_presentation.save(str(output_file))
                created_files.append(str(output_file))
                
                # Report progress - processing thumbnail (already generated)
                if progress_callback:
                    progress_callback(i, total_slides, slide_name, "creating_thumbnail")
                
                # Use pre-generated thumbnail from bulk conversion
                temp_thumbnail_path = bulk_thumbnail_paths[i-1] if i-1 < len(bulk_thumbnail_paths) else None
                
                # Resize and save the final thumbnail
                if temp_thumbnail_path:
                    thumbnail_path = self._process_and_save_thumbnail(temp_thumbnail_path, file_uuid)
                    # Clean up the temporary thumbnail
                    self.thumbnail_generator.cleanup_temp_thumbnail(temp_thumbnail_path)
                else:
                    # Fallback to individual generation if bulk failed for this slide
                    if self.config['verbose']:
                        print(f"    ⚠️  Using fallback thumbnail generation for slide {i}")
                    temp_thumbnail_path = self.thumbnail_generator.create_high_quality_thumbnail_from_pptx(
                        str(output_file), i
                    )
                    if temp_thumbnail_path:
                        thumbnail_path = self._process_and_save_thumbnail(temp_thumbnail_path, file_uuid)
                        self.thumbnail_generator.cleanup_temp_thumbnail(temp_thumbnail_path)
                    else:
                        thumbnail_path = None
                
                # Report progress - slide completed
                if progress_callback:
                    progress_callback(i, total_slides, slide_name, "completed")
                
                # Store slide metadata
                slide_metadata.append({
                    'name': slide_name,
                    'id': file_uuid,
                    'thumbMode': '1'
                })
                
                if self.config['verbose']:
                    print(f"✅ {output_file.name} + {Path(thumbnail_path).name}")
            
            # Generate XML metadata
            if progress_callback:
                progress_callback(total_slides + 1, total_slides + 2, "XML Metadata", "creating_xml")
            
            xml_path = self.xml_generator.create_xml_metadata(
                self.group_name, slide_metadata, self.output_dir
            )
            
            if self.config['verbose']:
                print(f"\n📄 Created XML metadata: {xml_path.name}")
                print(f"   Group: {self.group_name}")
                print(f"   Elements: {len(slide_metadata)}")
            
            # Create zip archive
            if progress_callback:
                progress_callback(total_slides + 2, total_slides + 2, "Zip Archive", "creating_zip")
            
            if self.config['verbose']:
                print(f"\n📦 Creating zip archive...")
            
            zip_path = self._create_zip_archive()
            
            # Final completion
            if progress_callback:
                progress_callback(total_slides + 2, total_slides + 2, "Export Complete", "export_complete")
            
            # Performance summary
            total_time = time.time() - start_time
            if self.config['verbose']:
                print(f"\n🎉 Processing complete!")
                print(f"⏱️  Total time: {total_time:.1f}s ({total_time/total_slides:.1f}s per slide)")
                print(f"🚀 Performance: {total_slides/total_time:.1f} slides/second")
                print(f"📈 High-quality thumbnails with accurate visual representation!")
            
            return created_files
            
        except Exception as e:
            # Clean up on error if we created a temporary directory
            if self.temp_dir_created:
                try:
                    cleanup_directory(self.output_dir, remove_directory=True, verbose=False)
                except:
                    pass
            raise
    
    def _create_single_slide_presentation(self, source_presentation, target_slide, slide_index):
        """
        Create a new presentation containing only the specified slide.
        
        Args:
            source_presentation: The original presentation
            target_slide: The slide to extract
            slide_index: Index of the slide in the original presentation
        
        Returns:
            A new Presentation object with only the target slide
        """
        # Create a temporary copy of the original presentation
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
            temp_path = Path(temp_file.name)
        
        # Save the source presentation to temp file
        source_presentation.save(str(temp_path))
        
        # Load the temp presentation
        temp_presentation = Presentation(str(temp_path))
        
        # Remove all slides except the target slide
        slides_to_remove = []
        for i, slide in enumerate(temp_presentation.slides):
            if i != slide_index:
                slides_to_remove.append(i)
        
        # Remove slides in reverse order to avoid index shifting
        for slide_idx in reversed(slides_to_remove):
            rId = temp_presentation.slides._sldIdLst[slide_idx].rId
            temp_presentation.part.drop_rel(rId)
            del temp_presentation.slides._sldIdLst[slide_idx]
        
        # Clean up temp file
        try:
            temp_path.unlink()
        except:
            pass
        
        return temp_presentation
    
    def _extract_slide_name(self, slide, slide_number: int) -> str:
        """
        Extract a meaningful name from the slide.
        
        Args:
            slide: The slide object
            slide_number: The slide number (1-based)
            
        Returns:
            The extracted slide name
        """
        # Try to find a title in the slide
        for shape in slide.shapes:
            try:
                if hasattr(shape, 'text') and shape.text.strip():
                    text_content = shape.text.strip()
                    
                    # Check if this might be a title
                    is_likely_title = False
                    
                    # Check if it's a placeholder
                    try:
                        if hasattr(shape, 'placeholder_format'):
                            is_likely_title = True
                    except:
                        pass
                    
                    # Or if it's short text (likely a title)
                    if len(text_content) < 100:
                        is_likely_title = True
                    
                    if is_likely_title:
                        # Clean up the title
                        title = ' '.join(text_content.split())
                        if len(title) > 50:
                            title = title[:47] + "..."
                        return title
            except Exception:
                # Skip shapes that cause errors
                continue
        
        # Fallback to generic slide name
        return f"Slide {slide_number}"
    
    def _process_and_save_thumbnail(self, temp_thumbnail_path: str, file_uuid: str) -> str:
        """
        Process and save the final thumbnail from a temporary thumbnail file.
        
        Args:
            temp_thumbnail_path: Path to the temporary thumbnail file
            file_uuid: UUID for the output filename
            
        Returns:
            Path to the final thumbnail file
        """
        try:
            # Final thumbnail path
            final_thumbnail_path = self.output_dir / f"{file_uuid}.png"
            
            # Resize the thumbnail to the configured height
            resized_thumbnail_path = self.thumbnail_generator.resize_thumbnail(
                temp_thumbnail_path, 
                self.config['thumbnail_height']
            )
            
            # Move the resized thumbnail to its final location
            import shutil
            shutil.move(resized_thumbnail_path, str(final_thumbnail_path))
            
            return str(final_thumbnail_path)
            
        except Exception as e:
            if self.config['verbose']:
                print(f"Error processing thumbnail: {e}")
            return None
    
    def _create_zip_archive(self) -> str:
        """
        Create a zip archive containing all generated files, then clean up.
        Places zip file at the same level as the input PowerPoint file.
        
        Returns:
            Path to the created zip file
        """
        try:
            # Create timestamped zip filename
            zip_filename = generate_timestamped_filename(self.base_name, "zip")
            
            # Place zip file at the same level as the input file
            zip_path = self.input_file.parent / zip_filename
            
            # Get all generated files to include in zip
            files_to_zip = []
            
            # Add all PPTX files (individual slides)
            files_to_zip.extend(self.output_dir.glob("*.pptx"))
            
            # Add all PNG files (thumbnails) 
            files_to_zip.extend(self.output_dir.glob("*.png"))
            
            # Add XML metadata file
            xml_file = self.output_dir / self.config['xml_filename']
            if xml_file.exists():
                files_to_zip.append(xml_file)
            
            if not files_to_zip:
                if self.config['verbose']:
                    print(f"    ⚠️  No files found to archive")
                return None
            
            # Create zip archive
            success, file_count, archive_size_mb = create_zip_archive(
                files_to_zip, 
                zip_path, 
                exclude_system_files=True
            )
            
            if success:
                if self.config['verbose']:
                    print(f"    ✅ Compressed {file_count} files")
                    print(f"    📦 Archive size: {archive_size_mb:.1f} MB")
                    print(f"    📁 Saved to: {zip_path}")
                
                # Clean up generated files now that they're archived
                if self.config['verbose']:
                    print(f"\n🧹 Cleaning up generated files...")
                
                cleanup_count = self._cleanup_generated_files()
                
                if self.config['verbose']:
                    print(f"    ✅ Removed {cleanup_count} generated files")
                    print(f"    📦 Final output: {zip_filename}")
                
                return str(zip_path)
            else:
                raise Exception("Failed to create zip archive")
                
        except Exception as e:
            if self.config['verbose']:
                print(f"    ❌ Error creating zip archive: {e}")
            return None
    
    def _cleanup_generated_files(self) -> int:
        """
        Remove all generated files (PPTX, PNG, XML) and the temporary directory.
        
        Returns:
            Number of files removed
        """
        # Get all files to remove
        files_to_remove = []
        file_patterns = ["*.pptx", "*.png", "*.xml"]
        
        for pattern in file_patterns:
            files_to_remove.extend(self.output_dir.glob(pattern))
        
        # Remove files
        removed_count = cleanup_files(files_to_remove, verbose=self.config['verbose'])
        
        # Clean up directory
        cleanup_directory(
            self.output_dir, 
            remove_directory=self.temp_dir_created,
            verbose=self.config['verbose']
        )
        
        return removed_count
